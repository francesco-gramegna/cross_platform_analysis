"""
Build the final group-presentation PowerPoint for DS516 Group 4.
Covers Data part + RQ1 (presenter's portion).

Output: ../../group_presentation/DS516_Group4_RQ1_final.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT = Path(__file__).resolve().parent
OUT_DIR = Path("/Users/draco/Desktop/2026 Spring/DS516/group_presentation")
OUT_DIR.mkdir(exist_ok=True, parents=True)
OUT = OUT_DIR / "DS516_Group4_RQ1_final.pptx"

# Plot paths
PLOTS_EDA = ROOT / "5 - EDA engagement" / "plots"
PLOTS_VOL = ROOT / "6 - tiktok volatility" / "plots"
PLOTS_RQ1 = ROOT / "7 - statistical analysis" / "plots"
PLOTS_ML  = ROOT / "8 - ML prediction" / "plots"
PLOTS_TEAM = ROOT / "EDA" / "figures"

# Design palette
NAVY    = RGBColor(0x1A, 0x23, 0x7E)   # primary - header bar
NAVY_LT = RGBColor(0x39, 0x49, 0xAB)   # lighter navy - accents
CORAL   = RGBColor(0xE5, 0x39, 0x35)   # accent red - emphasis
TEAL    = RGBColor(0x00, 0x83, 0x8F)   # secondary accent
ORANGE  = RGBColor(0xF5, 0x7C, 0x00)   # alert / scope
GREEN   = RGBColor(0x2E, 0x7D, 0x32)   # positive / bridge to RQ
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG = RGBColor(0xFA, 0xFA, 0xFC)   # subtle background tint
GRAY_LT = RGBColor(0xEE, 0xEE, 0xF2)   # light gray
TEXT    = RGBColor(0x21, 0x21, 0x21)
MUTED   = RGBColor(0x60, 0x60, 0x70)
DARK_HEADER = NAVY  # backwards compat
GRAY_LIGHT  = GRAY_LT
ACCENT      = CORAL

# Section assignments (slide index → section label)
SECTIONS = {
    range(2, 7):  "PART 1 · DATA",       # slides 2-6
    range(7, 10): "PART 2 · EDA",         # slides 7-9
    range(10, 17):"PART 3 · RQ1",         # slides 10-16
}
def section_for(slide_idx):
    for r, label in SECTIONS.items():
        if slide_idx in r: return label
    return ""

prs = Presentation()
prs.slide_width  = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

# ============================================================
# Helpers
# ============================================================
def add_title_only_slide(with_chrome=True):
    """Use blank layout. Add header bar + footer for chrome slides."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    if with_chrome:
        # subtle background tint
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        bg.fill.solid(); bg.fill.fore_color.rgb = GRAY_BG
        bg.line.fill.background()
        # send to back not directly supported; subsequent shapes will draw on top
        # header bar at very top
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.85))
        bar.fill.solid(); bar.fill.fore_color.rgb = NAVY
        bar.line.fill.background()
        # accent stripe under header (thin coral line)
        stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.85), prs.slide_width, Inches(0.06))
        stripe.fill.solid(); stripe.fill.fore_color.rgb = CORAL
        stripe.line.fill.background()
    return slide

def set_title(slide, text, size=22, bold=True, top=0.15, left=0.55, width=12.0,
              color=None):
    """Title text inside the navy header bar."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.55))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color if color else WHITE
    return tb

def add_section_tag(slide, section_label):
    """Small section indicator in top-right of header bar."""
    tb = slide.shapes.add_textbox(Inches(10.5), Inches(0.22), Inches(2.7), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    run = p.add_run(); run.text = section_label
    run.font.size = Pt(10); run.font.bold = True
    run.font.color.rgb = RGBColor(0xC5, 0xCA, 0xE9)  # light navy

def add_footer_bar(slide, slide_num, total=16):
    """Footer with slide number + project tag."""
    # small navy accent at bottom-left
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.35), Inches(0.4), Inches(0.15))
    accent.fill.solid(); accent.fill.fore_color.rgb = CORAL
    accent.line.fill.background()
    # text
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(7.18), Inches(12.3), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"DS516 · Group 4 · Cross-Platform Engagement"
    run.font.size = Pt(9); run.font.color.rgb = MUTED
    tb2 = slide.shapes.add_textbox(Inches(0.5), Inches(7.18), Inches(12.3), Inches(0.3))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    run2 = p2.add_run()
    run2.text = f"{slide_num} / {total}"
    run2.font.size = Pt(9); run2.font.bold = True; run2.font.color.rgb = NAVY

def add_subtitle(slide, text, top=1.0, left=0.5, width=12.3, size=14, color=None):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.italic = True
    run.font.color.rgb = color if color else RGBColor(0x60, 0x60, 0x60)

def add_bullets(slide, bullets, top=1.5, left=0.5, width=12.3, height=5.0, size=14, line_spacing=1.2):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(bullets):
        # bullet item can be:
        # - string
        # - dict: {"text": "...", "size": 14, "bold": False, "indent": 0}
        if isinstance(item, str):
            text, size_, bold_, indent_, color_ = item, size, False, 0, TEXT
        else:
            text = item.get("text","")
            size_ = item.get("size", size)
            bold_ = item.get("bold", False)
            indent_ = item.get("indent", 0)
            color_ = item.get("color", TEXT)
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.level = indent_
        run = p.add_run()
        if indent_ == 0 and not text.startswith("→") and not text.startswith("✦"):
            run.text = "• " + text
        else:
            run.text = text
        run.font.size = Pt(size_)
        run.font.bold = bold_
        run.font.color.rgb = color_

def add_image(slide, img_path, left, top, width=None, height=None):
    if not Path(img_path).exists():
        # placeholder text if missing
        tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(6), Inches(1))
        tf = tb.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"[Image not found: {Path(img_path).name}]"
        run.font.italic = True; run.font.color.rgb = RGBColor(0xCC, 0x00, 0x00)
        return
    kwargs = {}
    if width: kwargs["width"] = Inches(width)
    if height: kwargs["height"] = Inches(height)
    slide.shapes.add_picture(str(img_path), Inches(left), Inches(top), **kwargs)

def add_table(slide, data, left, top, width, height, header=True,
              col_widths=None, font_size=10):
    """data = list of rows, each row is a list of strings."""
    rows = len(data); cols = len(data[0])
    table_shape = slide.shapes.add_table(rows, cols, Inches(left), Inches(top),
                                          Inches(width), Inches(height))
    table = table_shape.table
    if col_widths:
        for j, w in enumerate(col_widths):
            table.columns[j].width = Inches(w)
    for i, row in enumerate(data):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = str(val)
            run.font.size = Pt(font_size)
            if i == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_HEADER
                run.font.color.rgb = WHITE
                run.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = GRAY_LIGHT if i % 2 == 0 else WHITE
                run.font.color.rgb = TEXT
            p.alignment = PP_ALIGN.CENTER

def add_footer(slide, text):
    tb = slide.shapes.add_textbox(Inches(0.4), Inches(7.05), Inches(12.5), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    run = p.add_run(); run.text = text
    run.font.size = Pt(9); run.font.italic = True
    run.font.color.rgb = RGBColor(0x80, 0x80, 0x80)

# ============================================================
# SLIDE 1 — Title (no chrome; full-bleed designed slide)
# ============================================================
s = add_title_only_slide(with_chrome=False)
# Full navy background
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
# Decorative coral block on left
block = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(0.3), Inches(2.5))
block.fill.solid(); block.fill.fore_color.rgb = CORAL; block.line.fill.background()
# Title text
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.5), Inches(11.5), Inches(2))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "Cross-Platform Influencer Engagement"
r.font.size = Pt(42); r.font.bold = True; r.font.color.rgb = WHITE
p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "What Shapes Online Audience Behavior?"
r.font.size = Pt(24); r.font.color.rgb = RGBColor(0xC5,0xCA,0xE9); r.font.italic = True
# Bottom info
tb = s.shapes.add_textbox(Inches(0.9), Inches(5.6), Inches(11.5), Inches(1.0))
tf = tb.text_frame
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "DS516 Final Presentation"
r.font.size = Pt(16); r.font.color.rgb = WHITE; r.font.bold = True
p = tf.add_paragraph(); p.alignment = PP_ALIGN.LEFT
r = p.add_run(); r.text = "Group 4 · [Member names]"
r.font.size = Pt(13); r.font.color.rgb = RGBColor(0xC5,0xCA,0xE9); r.font.italic = True
# Coral footer accent
foot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.2), prs.slide_width, Inches(0.3))
foot.fill.solid(); foot.fill.fore_color.rgb = CORAL; foot.line.fill.background()

# ============================================================
# SLIDE 2 — Motivation
# ============================================================
s = add_title_only_slide()
set_title(s, "What if engagement rate isn't a single thing?")
add_bullets(s, [
    {"text":"Three observations from our cleaned dataset:", "size":15, "bold":True},
    {"text":"TikTok creators earn ~8× the engagement rate of YouTube creators per follower (10.84% vs 1.37%)", "size":14},
    {"text":"On Instagram, doubling a creator's audience cuts engagement rate roughly in half (slope β ≈ −1)", "size":14},
    {"text":"On YouTube, the same doubling does almost nothing (slope β ≈ −0.3)", "size":14},
    {"text":" ", "size":8},
    {"text":"These aren't quirks — they reflect three fundamentally different platforms, each with its own interaction logic.", "size":14, "bold":True, "color":ACCENT},
    {"text":" ", "size":8},
    {"text":"Our question: What features actually drive engagement on each platform? Are the answers the same? What does it mean if they're not?", "size":14},
    {"text":" ", "size":8},
    {"text":"Scope: top 1,000 creators on IG, TT, YT across 2022 (5 monthly snapshots), 2024 (per-country leaderboards), 2026 (YT re-scrape).", "size":13, "color":RGBColor(0x60,0x60,0x60)},
], top=1.2)

# ============================================================
# SLIDE 3 — Data Sources (bullet layout)
# ============================================================
s = add_title_only_slide()
set_title(s, "Original Data Source")
add_bullets(s, [
    {"text":"Three sources from HypeAuditor (commercial influencer-analytics platform)", "size":15, "bold":True},
    {"text":" ", "size":6},
    {"text":"2022 — Kaggle HypeAuditor dump", "size":15, "bold":True, "color":ACCENT},
    {"text":"All three platforms (IG, TT, YT)", "size":13, "indent":1},
    {"text":"5 monthly snapshots (March / June / September / November / December)", "size":13, "indent":1},
    {"text":"Global top-1,000 per platform · 14,896 rows total", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"2024 — HypeAuditor per-country leaderboards", "size":15, "bold":True, "color":ACCENT},
    {"text":"All three platforms (IG, TT, YT)", "size":13, "indent":1},
    {"text":"Single snapshot · Top-100 per country across ~60 countries", "size":13, "indent":1},
    {"text":"17,280 rows total", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"2026 — HypeAuditor follow-up scrape", "size":15, "bold":True, "color":ACCENT},
    {"text":"YouTube only — revisits the 2022 cohort 4 years later", "size":13, "indent":1},
    {"text":"998 rows · ~45% are failed scrapes (deleted / inaccessible channels)", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"Key variables: followers · engagement count · category · audience country · platform · year (+ month for 2022)", "size":12, "color":RGBColor(0x60,0x60,0x60)},
], top=1.2)

# ============================================================
# SLIDE 4 — Preprocessing Pipeline (bullet layout)
# ============================================================
s = add_title_only_slide()
set_title(s, "Preprocessing Pipeline")
add_bullets(s, [
    {"text":"Five sequential cleaning steps applied across all three years", "size":15, "bold":True},
    {"text":" ", "size":6},
    {"text":"1.  Category unification", "size":14, "bold":True},
    {"text":"100+ raw category labels mapped to 8 standardized categories", "size":13, "indent":1},
    {"text":"2.  Category population", "size":14, "bold":True},
    {"text":"Fill missing categories using source's fallback fields (category_primary)", "size":13, "indent":1},
    {"text":"3.  Country normalization", "size":14, "bold":True},
    {"text":"ISO codes and native scripts converted to standardized country names", "size":13, "indent":1},
    {"text":"4.  Handle entity resolution", "size":14, "bold":True},
    {"text":"Match accounts across monthly snapshots and across years", "size":13, "indent":1},
    {"text":"5.  Engagement-rate unification", "size":14, "bold":True},
    {"text":"Build consistent engagement_count and er_pct using a single definition across sources", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"Methodological references", "size":13, "bold":True, "color":RGBColor(0x60,0x60,0x60)},
    {"text":"Entity resolution: Christen (2012), Data Matching", "size":11, "indent":1, "color":RGBColor(0x60,0x60,0x60)},
    {"text":"Category mapping: Euzenat & Shvaiko (2013), ontology alignment", "size":11, "indent":1, "color":RGBColor(0x60,0x60,0x60)},
    {"text":"ER definition: (likes + comments + shares) / followers — HypeAuditor, Hootsuite, Sprout Social standard", "size":11, "indent":1, "color":RGBColor(0x60,0x60,0x60)},
], top=1.0)

# ============================================================
# SLIDE 5 — Data Quality (bullet layout, by slice category)
# ============================================================
s = add_title_only_slide()
set_title(s, "Data Quality: What's Actually Usable?")
add_bullets(s, [
    {"text":"After cleaning, slices fall into three usability tiers:", "size":15, "bold":True},
    {"text":" ", "size":6},
    {"text":"✅ Fully usable — Instagram 2022 & YouTube 2022", "size":14, "bold":True, "color":RGBColor(0x1B,0x5E,0x20)},
    {"text":"ER coverage > 99%, category & audience country complete, no measurement artefacts", "size":12, "indent":1},
    {"text":" ", "size":6},
    {"text":"⚠️ Partial — TikTok 2022, all 2024 slices, YouTube 2026", "size":14, "bold":True, "color":RGBColor(0xE6,0x7E,0x22)},
    {"text":"TikTok 2022: no category / no audience country data at all", "size":12, "indent":1},
    {"text":"2024 (all platforms): source rounds 'er' to 0.1% → discretization artefacts in TT / YT", "size":12, "indent":1},
    {"text":"TikTok 2024 / YouTube 2024: low ER coverage (26% / 38%)", "size":12, "indent":1},
    {"text":"YouTube 2024: extremely sparse category coverage (6%)", "size":12, "indent":1},
    {"text":"YouTube 2026: ~45% are failed scrapes with no usable metrics", "size":12, "indent":1},
    {"text":" ", "size":6},
    {"text":"Sampling design note: 2022 is global top-1,000 (popularity-weighted); 2024 is top-100 per country (stratified). 'country' column means audience country in both — just sampled differently.", "size":12, "color":RGBColor(0x60,0x60,0x60)},
    {"text":" ", "size":6},
    {"text":"Scope decision: RQ1 inferential analysis restricted to 2022 — the only year with clean, popularity-weighted cross-platform data.", "size":13, "bold":True, "color":ACCENT},
], top=1.0)

# ============================================================
# SLIDE 6 — Research Questions (numbered card layout)
# ============================================================
s = add_title_only_slide()
set_title(s, "Research Questions")
add_bullets(s, [
    {"text":"Our team tackles engagement from four complementary angles:", "size":15, "bold":True},
    {"text":" ", "size":6},
    {"text":"RQ1 — What features drive engagement rate?", "size":15, "bold":True, "color":ACCENT},
    {"text":"Methods: OLS + ML.  Lead: [You]", "size":12, "indent":1, "color":RGBColor(0x60,0x60,0x60)},
    {"text":" ", "size":6},
    {"text":"RQ2 — Do platforms serve different content niches?", "size":15, "bold":True, "color":ACCENT},
    {"text":"Methods: [teammate's method].  Lead: [Teammate A]", "size":12, "indent":1, "color":RGBColor(0x60,0x60,0x60)},
    {"text":" ", "size":6},
    {"text":"RQ3 — Is the 'micro-influencer advantage' real, and does it vary by platform?", "size":15, "bold":True, "color":ACCENT},
    {"text":"Methods: [teammate's method].  Lead: [Teammate B]", "size":12, "indent":1, "color":RGBColor(0x60,0x60,0x60)},
    {"text":" ", "size":6},
    {"text":"RQ4 — Did content categories become more engagement-similar after AI tools?", "size":15, "bold":True, "color":ACCENT},
    {"text":"Methods: [teammate's method].  Lead: [Teammate C]", "size":12, "indent":1, "color":RGBColor(0x60,0x60,0x60)},
    {"text":" ", "size":6},
    {"text":"Common thread: how engagement structure differs across platforms, audience size, content type, or time.", "size":13, "color":RGBColor(0x60,0x60,0x60)},
    {"text":"Cross-RQ flow: EDA sets the stage → RQ1 identifies drivers → RQ2/3/4 dig into specific dimensions.", "size":13, "color":RGBColor(0x60,0x60,0x60)},
], top=1.0)

# ============================================================
# SLIDE 7 — EDA 1/3: ER hierarchy → RQ1
# ============================================================
s = add_title_only_slide()
set_title(s, "Platforms differ massively in engagement")
add_image(s, PLOTS_EDA / "P3_er_pct_boxplot_2022.png", left=0.5, top=1.2, width=7.5)
add_bullets(s, [
    {"text":"At a glance:", "size":14, "bold":True},
    {"text":"TikTok median ER: 10.8%", "size":14},
    {"text":"Instagram median ER: 3.5%", "size":14},
    {"text":"YouTube median ER: 1.4%", "size":14},
    {"text":" ", "size":6},
    {"text":"TikTok creators earn roughly 8× YouTube's engagement per follower.", "size":13, "bold":True, "color":ACCENT},
    {"text":" ", "size":6},
    {"text":"These numbers come from the same kind of accounts (global top-1,000 in 2022), yet platforms aren't even close.", "size":13},
    {"text":" ", "size":6},
    {"text":"→ What features drive engagement on each platform? (RQ1)", "size":14, "bold":True, "color":RGBColor(0x1B,0x5E,0x20)},
], top=1.4, left=8.3, width=4.6)

# ============================================================
# SLIDE 8 — EDA 2/3: Followers vs ER → RQ3
# ============================================================
s = add_title_only_slide()
set_title(s, "Smaller accounts seem to engage more")
add_image(s, PLOTS_EDA / "P11_followers_vs_er_clean_scatter.png", left=0.5, top=1.2, width=7.8)
add_bullets(s, [
    {"text":"At a glance:", "size":14, "bold":True},
    {"text":"On every platform, the cloud slopes downward — bigger accounts engage less per follower", "size":13},
    {"text":"Trend is steep on Instagram, very pronounced on TikTok, gentle on YouTube", "size":13},
    {"text":" ", "size":6},
    {"text":"TikTok is the only platform where our top-1,000 sample contains small accounts — its leaderboard churns fast (~41% of TT accounts appear in only one of 5 monthly snapshots).", "size":12, "color":RGBColor(0x60,0x60,0x60)},
    {"text":" ", "size":6},
    {"text":"→ Is the 'micro-influencer advantage' real and uniform across platforms? (RQ3)", "size":14, "bold":True, "color":RGBColor(0x1B,0x5E,0x20)},
], top=1.4, left=8.6, width=4.3)

# ============================================================
# SLIDE 9 — EDA 3/3: Content niches → RQ2 + RQ4
# ============================================================
s = add_title_only_slide()
set_title(s, "Platforms specialize, and content is shifting over time")
# Use team plot
team_plot = PLOTS_TEAM / "cat_distribution_by_platform_year.png"
if team_plot.exists():
    add_image(s, team_plot, left=0.5, top=1.2, width=7.5)
else:
    add_image(s, PLOTS_EDA / "P9_category_niche_by_platform.png", left=0.5, top=1.2, width=7.5)

add_bullets(s, [
    {"text":"Music + Entertainment dominate every platform — but each has over-represented categories:", "size":13, "bold":True},
    {"text":"Tech & Gaming + Knowledge & Info → YouTube (~10% each, vs 1-3% on IG/TT)", "size":12},
    {"text":"Beauty & Fashion → TikTok (15-18%, vs 1% on YT)", "size":12},
    {"text":"Sports → Instagram (14%, vs 1% on YT)", "size":12},
    {"text":" ", "size":6},
    {"text":"And content mix is shifting over time:", "size":13, "bold":True},
    {"text":"IG Beauty & Fashion: 8% → 29% (2022→2024)", "size":12},
    {"text":"IG Entertainment: 28% → 9% (2022→2024)", "size":12},
    {"text":" ", "size":6},
    {"text":"→ Are these niches structural? (RQ2)", "size":13, "bold":True, "color":RGBColor(0x1B,0x5E,0x20)},
    {"text":"→ Did AI tools push categories to converge? (RQ4)", "size":13, "bold":True, "color":RGBColor(0x1B,0x5E,0x20)},
], top=1.3, left=8.3, width=4.6)

# ============================================================
# SLIDE 10 — RQ1: setup
# ============================================================
s = add_title_only_slide()
set_title(s, "RQ1: What Drives Engagement Rate?")
add_bullets(s, [
    {"text":"Question: which account features predict ER, and do effects look the same across platforms?", "size":15, "bold":True},
    {"text":" ", "size":6},
    {"text":"Target: log10(er_pct) — log handles right-skew, satisfies OLS assumptions", "size":14},
    {"text":"Predictors: log10(followers), content category, audience country (top-15)", "size":14},
    {"text":"References: Entertainment (category), United States (country)", "size":14},
    {"text":"Scope: 2022, Instagram + YouTube (TikTok lacks usable category & country)", "size":14},
    {"text":" ", "size":6},
    {"text":"Two-stage analysis:", "size":14, "bold":True},
    {"text":"Stage 1 (Traditional): Pearson · ANOVA · multiple OLS with interactions · Welch t-test", "size":13},
    {"text":"Stage 2 (Advanced): OLS baseline · Random Forest · XGBoost · SHAP", "size":13},
    {"text":" ", "size":6},
    {"text":"Lit anchor: De Veirman, Cauberghe & Hudders (2017) — follower count negatively predicts ER on Instagram. We test whether this pattern generalizes.", "size":12, "color":RGBColor(0x60,0x60,0x60)},
], top=1.2)

# ============================================================
# SLIDE 11 — RQ1 Methods at a Glance (bullet layout)
# ============================================================
s = add_title_only_slide()
set_title(s, "RQ1: Methods at a Glance")
add_bullets(s, [
    {"text":"Four statistical methods, each answering a different question:", "size":15, "bold":True},
    {"text":" ", "size":6},
    {"text":"Pearson r — Is followers ↔ ER associated?", "size":14, "bold":True},
    {"text":"Result: negative on every platform; Instagram strongest (r = −0.73)", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"ANOVA on category — Does content category matter?", "size":14, "bold":True},
    {"text":"Result: significant on IG (η² = 0.08) and YT (η² = 0.13); effect is stronger on YouTube", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"Multiple OLS with interactions — What is each feature's controlled effect?", "size":14, "bold":True},
    {"text":"Result: Instagram β_followers = −0.99, R² = 0.57.  YouTube β_followers = −0.34, R² = 0.22", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"Welch t-test — Do ER levels differ across platforms?", "size":14, "bold":True},
    {"text":"Result: TikTok > Instagram > YouTube, all pairs p ≈ 0", "size":13, "indent":1},
    {"text":" ", "size":6},
    {"text":"Headline: followers is the universal driver; category and country effects are platform-specific in magnitude AND sign.", "size":13, "bold":True, "color":ACCENT},
    {"text":"The numbers are the foundation — the next slides focus on what they actually mean.", "size":12, "color":RGBColor(0x60,0x60,0x60)},
], top=1.0)

# ============================================================
# SLIDE 12 — RQ1 OLS vs ML
# ============================================================
s = add_title_only_slide()
set_title(s, "RQ1 (Advanced): Does ML beat the linear baseline?")
add_image(s, PLOTS_ML / "F1_predicted_vs_actual.png", left=0.4, top=1.2, width=8.0)
add_bullets(s, [
    {"text":"Three models, 80/20 train/test split, same target & features", "size":14, "bold":True},
    {"text":" ", "size":6},
    {"text":"OLS (baseline): R² = 0.229", "size":13},
    {"text":"Random Forest: R² = 0.146 (overfits)", "size":13},
    {"text":"XGBoost: R² = 0.223 (ties OLS)", "size":13},
    {"text":" ", "size":6},
    {"text":"ML does NOT outperform OLS → predictable signal is mostly linear.", "size":13, "bold":True, "color":ACCENT},
    {"text":" ", "size":6},
    {"text":"All three models agree on ranking: followers > category > country > platform (permutation + SHAP)", "size":12},
    {"text":" ", "size":6},
    {"text":"R² ≈ 0.23 is the realistic ceiling. The other 77% lives in unobserved post-level signals.", "size":12, "color":RGBColor(0x60,0x60,0x60)},
], top=1.3, left=8.6, width=4.3)

# ============================================================
# SLIDE 13 — Finding 1: Followers is the WHOLE story on IG, barely matters on YT
# ============================================================
s = add_title_only_slide()
set_title(s, "Finding 1: Followers is everything on Instagram — barely matters on YouTube")
add_image(s, PLOTS_RQ1 / "F10_followers_slope_contrast.png", left=0.3, top=1.0, width=7.6)
add_bullets(s, [
    {"text":"SETUP", "size":12, "bold":True, "color":NAVY},
    {"text":"How much does follower count predict engagement on each platform?", "size":11},
    {"text":" ", "size":6},
    {"text":"EVIDENCE 1 — Followers explains 80× MORE variance on Instagram", "size":12, "bold":True, "color":NAVY},
    {"text":"Followers-only ΔR²:  IG = 47.6%   vs   YT = 0.6%", "size":11, "bold":True},
    {"text":"On IG, follower count is the SINGLE largest predictor of ER", "size":11},
    {"text":"On YT, follower count barely matters", "size":11},
    {"text":" ", "size":6},
    {"text":"EVIDENCE 2 — Slope contrast confirms the variance gap", "size":12, "bold":True, "color":NAVY},
    {"text":"Instagram β = −0.99 → 10× more followers ≈ 10× lower ER", "size":11},
    {"text":"YouTube β = −0.34 → 10× more followers ≈ 2× lower ER", "size":11},
    {"text":"Interaction p < 10⁻⁹ — slopes are statistically different", "size":11},
    {"text":" ", "size":6},
    {"text":"WHY", "size":12, "bold":True, "color":NAVY},
    {"text":"Instagram dilution: as accounts grow, the engaged-follower share collapses (millions of casual fans dilute the original tight community)", "size":11},
    {"text":"YouTube structural passivity: subscribers consume passively at every account size; viral views drive ER more than audience size", "size":11},
    {"text":" ", "size":6},
    {"text":"Takeaway: on Instagram, follower size IS the engagement story. On YouTube, you need to look at WHAT content you make and WHO watches it (Findings 2 & 3).", "size":11, "bold":True, "color":CORAL},
], top=1.0, left=7.9, width=5.1)

# ============================================================
# SLIDE 14 — Finding 2: Category defines engagement on YT, not on IG
# ============================================================
s = add_title_only_slide()
set_title(s, "Finding 2: Content category defines engagement on YouTube — not Instagram")
add_image(s, PLOTS_RQ1 / "F11_category_general_pattern.png", left=0.2, top=1.0, width=7.6)
add_bullets(s, [
    {"text":"SETUP", "size":12, "bold":True, "color":NAVY},
    {"text":"Does content category drive ER the same way on each platform?", "size":11},
    {"text":" ", "size":6},
    {"text":"EVIDENCE 1 — Category contributes 12× MORE variance on YouTube", "size":12, "bold":True, "color":NAVY},
    {"text":"Category ΔR²:  IG +1.2%   vs   YT +14.6%", "size":11},
    {"text":"On YT, category is one of the LARGEST predictors of ER", "size":11, "bold":True},
    {"text":" ", "size":6},
    {"text":"EVIDENCE 2 — On Instagram, ALL bars are tiny (~0)", "size":12, "bold":True, "color":NAVY},
    {"text":"Every category sits near the Entertainment baseline", "size":11},
    {"text":"→ category barely shifts engagement on IG", "size":11, "color":MUTED},
    {"text":" ", "size":6},
    {"text":"EVIDENCE 3 — On YouTube, bars are MUCH larger and in BOTH directions", "size":12, "bold":True, "color":NAVY},
    {"text":"Strong negative on YT (categories that engage WORSE than Entertainment):", "size":11, "bold":True},
    {"text":"  Knowledge & Info: β = −0.76", "size":11},
    {"text":"  Sports: β = −0.56", "size":11},
    {"text":"  Music: β = −0.35", "size":11},
    {"text":"Strong positive on YT (engages BETTER):", "size":11, "bold":True},
    {"text":"  Tech & Gaming: β = +0.18 (the only positive lift)", "size":11},
    {"text":" ", "size":6},
    {"text":"WHY", "size":12, "bold":True, "color":NAVY},
    {"text":"YT cultivates niche SUBCULTURES — each content type has its own audience interaction style (gamers debate, learners just watch)", "size":11},
    {"text":"IG is UNIFORM-BROWSING — visual feed treats all categories alike, so category type barely shifts engagement", "size":11},
    {"text":" ", "size":6},
    {"text":"Takeaway: on YouTube, content niche IS destiny. On Instagram, follower dilution dominates and category is noise.", "size":11, "bold":True, "color":CORAL},
], top=1.0, left=7.9, width=5.1)

# ============================================================
# SLIDE 15 — Finding 3: Country effect has TWO distinct mechanisms
# ============================================================
s = add_title_only_slide()
set_title(s, "Finding 3: The 'country effect' has two distinct mechanisms")
add_image(s, PLOTS_RQ1 / "F19_country_by_region.png", left=0.2, top=1.0, width=7.6)
add_bullets(s, [
    {"text":"SETUP", "size":12, "bold":True, "color":NAVY},
    {"text":"Does audience country drive ER the same way on Instagram and YouTube?", "size":11},
    {"text":" ", "size":6},
    {"text":"EVIDENCE 1 — Country contributes 3× MORE variance on YouTube", "size":12, "bold":True, "color":NAVY},
    {"text":"Country ΔR²:  IG +1.9%  vs  YT +6.3%", "size":11},
    {"text":"Direction flips: IG 3↑/0↓ vs YT 1↑/5↓ (sig at p<0.05)", "size":11},
    {"text":" ", "size":6},
    {"text":"EVIDENCE 2 — Engagement shifts go in OPPOSITE directions by region", "size":12, "bold":True, "color":NAVY},
    {"text":"⬆ MORE on YT: Eastern Europe (3.6% → 6.4%), W. Europe (3.7% → 4.4%)", "size":11, "color":GREEN},
    {"text":"⬇ MUCH LESS on YT: South Asia (2.7% → 0.4%), SE Asia (6.1% → 1.0%)", "size":11, "color":CORAL},
    {"text":"Moderate decline: North America, Latin America, MENA", "size":11},
    {"text":" ", "size":6},
    {"text":"WHY — two different mechanisms behind the drops", "size":12, "bold":True, "color":NAVY},
    {"text":" ", "size":6},
    {"text":"South Asia = BEHAVIORAL difference", "size":12, "bold":True, "color":CORAL},
    {"text":"Within each category, Indian audiences engage less per follower than US", "size":11},
    {"text":"e.g., Indian K&I creators = 0.11% vs US K&I = 1.86% (17× lower)", "size":11},
    {"text":"→ Same content type, but audience genuinely engages less", "size":11, "color":MUTED},
    {"text":" ", "size":6},
    {"text":"SE Asia = CONTENT-MIX difference", "size":12, "bold":True, "color":CORAL},
    {"text":"Within each category, Indonesia roughly matches US engagement", "size":11},
    {"text":"e.g., Indonesian Music = 1.48% vs US Music = 1.07% (HIGHER!)", "size":11},
    {"text":"→ The aggregate gap reflects which categories the region has, not how its audiences engage", "size":11, "color":MUTED},
    {"text":" ", "size":6},
    {"text":"Takeaway: a single country coefficient hides this — the same 'negative country effect' can be behavioral in one region and just content-mix in another.", "size":11, "bold":True, "color":CORAL},
], top=1.0, left=7.9, width=5.1)

# ============================================================
# SLIDE 16 — RQ1 Limitations
# ============================================================
s = add_title_only_slide()
set_title(s, "RQ1: Limitations")
add_bullets(s, [
    {"text":"1.  Data scope", "size":13, "bold":True, "color":NAVY},
    {"text":"2022 only — no temporal trends in RQ1's main scope", "size":11, "indent":1},
    {"text":"Top-1,000 creators per platform — survivorship bias toward popularity", "size":11, "indent":1},
    {"text":"TikTok 2022 excluded — no category or audience-country data available", "size":11, "indent":1},
    {"text":"Some regions sparse: East Asia (S. Korea absent on YT), MENA, Africa & Oceania absent entirely", "size":11, "indent":1},
    {"text":" ", "size":6},
    {"text":"2.  Measurement comparability across platforms", "size":13, "bold":True, "color":NAVY},
    {"text":"Engagement-count defined differently per platform:", "size":11, "indent":1},
    {"text":"  IG = HypeAuditor's proprietary engagement_total (opaque formula)", "size":11, "indent":1},
    {"text":"  TT = likes + comments + shares", "size":11, "indent":1},
    {"text":"  YT = likes + comments + views (treats deliberate click as engagement)", "size":11, "indent":1},
    {"text":"Cross-platform ER multipliers should be read as ORDINAL, not precise", "size":11, "indent":1},
    {"text":"HypeAuditor's curation choices are not fully transparent", "size":11, "indent":1},
    {"text":" ", "size":6},
    {"text":"3.  Statistical methodology", "size":13, "bold":True, "color":NAVY},
    {"text":"Category × Country interaction NOT tested — too many cells empty / underpowered", "size":11, "indent":1},
    {"text":"Variance decomposition is sequential — order of variable entry affects ΔR²", "size":11, "indent":1},
    {"text":"Within-category cells have small n in some places (visually faded)", "size":11, "indent":1},
    {"text":" ", "size":6},
    {"text":"4.  Causal vs correlational", "size":13, "bold":True, "color":NAVY},
    {"text":"All findings are observational — no causal claims", "size":11, "indent":1},
    {"text":"Mechanism interpretations (gaming for E. Europe; behavioral for S. Asia) are supported but NOT directly tested", "size":11, "indent":1},
    {"text":"Unobserved confounders likely: creator language, posting frequency, content quality, finer-grained audience demographics", "size":11, "indent":1},
], top=1.0)

# ============================================================
# SLIDE 17 — Synthesis
# ============================================================
s = add_title_only_slide()
set_title(s, "Synthesis: What RQ1 Tells Us — and What's Next")
add_bullets(s, [
    {"text":"Three findings worth taking forward:", "size":16, "bold":True},
    {"text":"Followers effect differs by platform — IG slope ≈ 3× steeper than YT", "size":13},
    {"text":"Category effects flip across platforms — Lifestyle leads IG; Tech & Gaming leads YT", "size":13},
    {"text":"Country direction flips across platforms — non-US engages more on IG, less on YT", "size":13},
    {"text":" ", "size":6},
    {"text":"Methodological note:", "size":14, "bold":True},
    {"text":"OLS / RF / XGBoost converge at R² ≈ 0.23 → signal is mostly linear", "size":12},
    {"text":"~77% of ER variance lives in unobserved post-level signals", "size":12},
    {"text":" ", "size":6},
    {"text":"What this raises for the rest of the presentation:", "size":14, "bold":True, "color":ACCENT},
    {"text":"Platforms have different category drivers → Do they actually serve different content niches? (RQ2)", "size":12},
    {"text":"Followers strongly predict ER, with platform-specific strength → Does the micro-influencer advantage hold uniformly? (RQ3)", "size":12},
    {"text":"Category effects can shift over time → Did the AI era make content categories more similar? (RQ4)", "size":12},
    {"text":" ", "size":6},
    {"text":"Methodological takeaway: engagement rate cannot be modeled as one platform-agnostic quantity. Per-platform analysis is required.", "size":13, "bold":True, "color":ACCENT},
    {"text":" ", "size":6},
    {"text":"→ Hand off to [teammate] for RQ2", "size":14, "bold":True, "color":RGBColor(0x1B,0x5E,0x20)},
], top=1.2)

# Add section tag + footer to all content slides (skip title slide #1)
total = len(prs.slides)
for idx, slide in enumerate(prs.slides, start=1):
    if idx == 1:  # title slide — already designed
        continue
    section = section_for(idx)
    if section:
        add_section_tag(slide, section)
    add_footer_bar(slide, idx, total=total)

# Save
prs.save(OUT)
print(f"Saved: {OUT}")
print(f"File size: {OUT.stat().st_size:,} bytes")
print(f"Slides: {len(prs.slides)}")
